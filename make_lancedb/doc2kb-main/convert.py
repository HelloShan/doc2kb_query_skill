"""
doc2kb — 文档转换引擎模块
=============================
转换引擎优先级：
  主力：Docling（支持 docx/pdf/xlsx/pptx，高质量表格还原+多栏识别）
  降级：原生 Python 解析器（python-docx / pypdf / openpyxl / python-pptx）

支持将 6 种核心格式转换为清洁后的 Markdown：
  - .docx  → Docling(主) / python-docx(降)
  - .pdf   → Docling(主) / pypdf(降)
  - .xlsx  → Docling(主) / openpyxl(降)
  - .pptx  → Docling(主) / python-pptx(降)
  - .md    → 复制+乱码校验
  - .txt   → 直接复制，编码自动回退

每个文件返回 (status, md_rel_path_or_None, error_msg_or_None, warning_or_None) 四元组。
"""

import os
import re
import sys
import threading
import traceback
from pathlib import Path
from typing import Optional, Tuple, List
from concurrent.futures import ThreadPoolExecutor, as_completed

from config import (
    SOURCE_DIR, OUTPUT_MD_DIR, SUPPORTED_EXTENSIONS, CODE_CONFIG_EXTENSIONS,
    SOURCE_EXT_MARKER_PREFIX,
    CONVERT_WORKERS, MAX_MD_FILE_SIZE_KB, CONVERT_TIMEOUT, DOCLING_MAX_CONCURRENT,
)
from validate import is_file_readable
from state import compute_sha256
from docx.opc.exceptions import PackageNotFoundError

# 模块级 /dev/null 句柄（多线程复用，避免并发 close 炸 sys.stderr）
_DEVNULL = open(os.devnull, "w")


# ============================================================
# 工具函数
# ============================================================

def _get_rel_path(source_path: Path) -> str:
    return source_path.relative_to(SOURCE_DIR).as_posix()


def _get_output_md_path(source_path: Path) -> Path:
    rel = source_path.relative_to(SOURCE_DIR)
    return OUTPUT_MD_DIR / rel.with_suffix(".md")


def _ensure_output_dir(output_path: Path):
    output_path.parent.mkdir(parents=True, exist_ok=True)


# ============================================================
# Docling 引擎（主力）
# ============================================================

_DOCLING_CONVERTER = None
_DOCLING_SUPPORTED = {".docx", ".pdf"}

# Docling 单独的并发闸门，独立于 CONVERT_WORKERS（见 config.py 里
# DOCLING_MAX_CONCURRENT 的详细说明）。docx/pdf 转换在拿到这个信号量之前
# 会排队等待，不会真正去起那个会加载一整套模型的子进程；xlsx/pptx/txt
# 等原生解析路径完全不受影响，仍按 CONVERT_WORKERS 的并发数正常跑。
_docling_semaphore = threading.Semaphore(max(1, DOCLING_MAX_CONCURRENT))


def _get_docling_converter():
    global _DOCLING_CONVERTER
    if _DOCLING_CONVERTER is None:
        from docling.document_converter import DocumentConverter, PdfFormatOption
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import PdfPipelineOptions
        from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend

        # 不需要 OCR：知识库源文档都是可直接抽取文本层的电子版
        # docx/pdf，不是扫描件。Docling 默认会加载 OCR 模型（RapidOCR）
        # 处理"疑似缺少文本层"的区域（比如页面里的截图、大图），这一步
        # 本来就用不上，显式关掉，识别不了文字的图片/扫描内容直接舍弃。
        pipeline_options = PdfPipelineOptions()
        pipeline_options.do_ocr = False

        # PDF 解析后端换成 pypdfium2，不用 Docling 默认的 docling-parse。
        # ────────────────────────────────────────────────────────────
        # docling-parse 这个 C++ 后端在解析包含大图片的 PDF 时，有已知
        # 的"preprocess 阶段按页累积内存直至 std::bad_alloc"问题（见
        # docling 官方仓库 issue #3345、#3671，docling-parse 仓库
        # issue #227 等，都是同一个报错特征：
        #   "Stage preprocess failed for run 1, pages [N]: std::bad_alloc"
        # ）。这份"硬件安装指导"文档图片密集，120 页左右单次转换到第
        # 117 页左右就撞上了这个问题——关掉 OCR 之后依然会崩，说明和
        # OCR 无关，是这个底层解析后端本身的缺陷。
        #
        # 多个 issue 里都确认换成 pypdfium2 后端可以完全规避这个问题
        # （代价是复杂表格的结构还原质量可能不如 docling-parse，但能
        # 稳定把文件转完，比转到一半直接把子进程崩掉、整份文件转换
        # 失败要好）。
        _DOCLING_CONVERTER = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(
                    pipeline_options=pipeline_options,
                    backend=PyPdfiumDocumentBackend,
                )
            }
        )
    return _DOCLING_CONVERTER


# ============================================================
# 兼容性检测（供 doc_pipeline.py check 子命令及转换前预检共用）
# ============================================================

def detect_docm(file_path: Path) -> tuple[bool, str]:
    """检测 .docx 文件是否为宏文档 (.docm) 或旧版 .doc。返回 (是问题吗, 原因)。"""
    if file_path.suffix.lower() not in ('.docx',):
        return False, ''
    try:
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(str(file_path)) as z:
            ct = z.read('[Content_Types].xml')
            root = ET.fromstring(ct)
            ns = root.tag.split('}')[0].strip('{') if '}' in root.tag else ''
            tag = f'{{{ns}}}Override' if ns else 'Override'
            for override in root.iter(tag):
                ct_type = override.get('ContentType', '')
                if 'macroenabled' in ct_type.lower():
                    return True, '宏文档 (.docm)，需用 Word 另存为 .docx'
        return False, ''
    except (zipfile.BadZipFile, Exception):
        return True, '不是有效的 ZIP 包（可能是旧版 .doc 格式误标为 .docx）'


def detect_legacy_xls(file_path: Path) -> tuple[bool, str]:
    """
    检测 .xlsx 文件是否其实是旧版 .xls（二进制 OLE 格式）误标了扩展名。
    .xlsx 本质是 ZIP 包，打不开说明大概率是老格式内容、新格式扩展名。
    返回 (是问题吗, 原因)。
    """
    if file_path.suffix.lower() != '.xlsx':
        return False, ''
    try:
        import zipfile
        with zipfile.ZipFile(str(file_path)):
            pass
        return False, ''
    except zipfile.BadZipFile:
        return True, '不是有效的 ZIP 包（可能是旧版 .xls 格式误标为 .xlsx）'
    except Exception as e:
        return True, f'读取失败（可能是旧版 .xls 格式误标为 .xlsx）: {type(e).__name__}'


# ============================================================
# 老版 Office 格式（误标扩展名）自动转换
# ============================================================
# 场景：源文件后缀写的是 .docx/.xlsx，但内容其实是老版二进制格式
# .doc/.xls（比如批量改后缀、或者从别处复制过来时手滑改错）。这类
# 文件 python-docx/openpyxl/Docling 都打不开，之前的做法是直接跳过，
# 提示"需要手工用 WPS/Word 打开后另存为新格式、删除老文件、重新转换"。
#
# 这里把这个手工步骤自动化：调用 WPS（优先，因为很多内部环境没装
# 正版 Office）或 Microsoft Office 的 COM 自动化，把文件原地另存为
# 真正的新格式，后续转换流程不用改，正常按 .docx/.xlsx 走 Docling。
#
# 仅支持 Windows（COM 是 Windows 专属机制）。非 Windows 环境，或者
# WPS/Office 都没装、pywin32 没装，会直接返回失败，调用方回退到原来
# "跳过 + 提示手工处理"的行为，不会导致其它平台/环境跑不起来。
_COM_CONVERT_LOCK = threading.Lock()


def _convert_legacy_office_file(source_path: Path, kind: str) -> Tuple[bool, str]:
    """
    把 source_path（扩展名是新格式，内容是老格式）原地转换为真正的新格式。

    Parameters
    ----------
    kind : "word" | "excel"

    Returns
    -------
    (成功与否, 说明信息)
    """
    if sys.platform != "win32":
        return False, "自动转换仅支持 Windows（需要 WPS 或 Office 的 COM 自动化）"

    try:
        import win32com.client
        import pythoncom
    except ImportError:
        return False, "缺少 pywin32，无法调用 WPS/Office COM 自动化（pip install pywin32）"

    import tempfile
    import shutil
    import zipfile

    legacy_ext = ".doc" if kind == "word" else ".xls"
    new_ext = ".docx" if kind == "word" else ".xlsx"

    tmp_dir = Path(tempfile.mkdtemp(prefix="doc2kb_legacy_"))
    # 先把内容复制成扩展名和内容匹配的临时文件（老格式扩展名），
    # 避免 WPS/Office 因为"扩展名和内容对不上"而拒绝打开或弹确认框。
    legacy_tmp = tmp_dir / (source_path.stem + legacy_ext)
    new_tmp = tmp_dir / (source_path.stem + new_ext)

    try:
        shutil.copy2(source_path, legacy_tmp)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        return False, f"复制临时文件失败: {type(e).__name__}: {e}"

    # COM 自动化不是线程安全的，这个流水线用线程池并行转换文件，这里
    # 用一把全局锁把"调 WPS/Office 转换"串行化。这类误标扩展名的文件
    # 本来就是极少数，串行化不会成为性能瓶颈，但能避免多个线程同时
    # 操作 COM 对象导致的各种诡异报错、卡死或残留僵尸进程。
    with _COM_CONVERT_LOCK:
        pythoncom.CoInitialize()
        app = None
        doc = None
        used_prog_id = None
        try:
            if kind == "word":
                # KWPS.Application = WPS 文字；Word.Application = MS Word。
                # 优先 WPS，因为很多内部环境只装了 WPS 没装正版 Office。
                prog_ids = ["KWPS.Application", "Word.Application"]
                wd_format_docx = 16  # wdFormatXMLDocument，WPS 兼容同一个值
            else:
                prog_ids = ["KET.Application", "Excel.Application"]
                xl_format_xlsx = 51  # xlOpenXMLWorkbook，WPS 兼容同一个值

            dispatch_errors = []
            for prog_id in prog_ids:
                try:
                    app = win32com.client.DispatchEx(prog_id)
                    used_prog_id = prog_id
                    break
                except Exception as e:
                    dispatch_errors.append(f"{prog_id}: {type(e).__name__}: {e}")
                    app = None

            if app is None:
                return False, ("WPS 和 Office 都无法通过 COM 启动（可能都没装，"
                                "或没有注册 COM 组件）: " + "; ".join(dispatch_errors))

            app.Visible = False
            try:
                app.DisplayAlerts = False
            except Exception:
                pass  # 个别版本没有这个属性也无所谓，不是必需的

            if kind == "word":
                doc = app.Documents.Open(str(legacy_tmp), ConfirmConversions=False,
                                          ReadOnly=False, AddToRecentFiles=False)
                doc.SaveAs2(str(new_tmp), FileFormat=wd_format_docx)
            else:
                doc = app.Workbooks.Open(str(legacy_tmp), UpdateLinks=0,
                                          ReadOnly=False, AddToMru=False)
                doc.SaveAs(str(new_tmp), FileFormat=xl_format_xlsx)

            doc.Close(False)
            doc = None

            if not new_tmp.exists() or new_tmp.stat().st_size == 0:
                return False, "另存为新格式后文件为空或不存在"

            # .docx/.xlsx 本质是 ZIP 包，校验一下另存出来的文件确实合法，
            # 避免把一个损坏的转换结果替换掉原文件。
            try:
                with zipfile.ZipFile(str(new_tmp)):
                    pass
            except zipfile.BadZipFile:
                return False, "另存为新格式后的文件不是合法的 ZIP 包，转换结果可能已损坏"

            # 原地替换：先把老格式内容备份一份，再用新格式覆盖原路径。
            # os.replace 是原子操作，不会出现"写到一半、文件半新半旧"的中间状态。
            backup_path = source_path.with_name(source_path.name + ".legacy_bak")
            shutil.copy2(source_path, backup_path)
            os.replace(str(new_tmp), str(source_path))

            return True, f"已用 {used_prog_id} 自动转换为新格式（原文件已备份为 {backup_path.name}）"

        except Exception as e:
            return False, f"COM 自动转换失败: {type(e).__name__}: {e}"
        finally:
            try:
                if doc is not None:
                    doc.Close(False)
            except Exception:
                pass
            try:
                if app is not None:
                    app.Quit()
            except Exception:
                pass
            pythoncom.CoUninitialize()
            shutil.rmtree(tmp_dir, ignore_errors=True)


def detect_broken_pdf(file_path: Path) -> tuple[bool, str]:
    """检查 PDF 文件是否可读（静默模式，压制 pypdf 的 stderr 和日志噪音）。"""
    import logging
    import contextlib

    if file_path.suffix.lower() != '.pdf':
        return False, ''
    try:
        from pypdf import PdfReader, errors as pypdf_errors
        # 静默 pypdf 的日志和 stderr 噪音（如 incorrect startxref pointer）
        logger = logging.getLogger('pypdf')
        old_level = logger.level
        logger.setLevel(logging.ERROR)
        try:
            with contextlib.redirect_stderr(_DEVNULL):
                reader = PdfReader(str(file_path))
                _ = len(reader.pages)
            return False, ''
        except pypdf_errors.PdfStreamError:
            return True, 'PDF 流意外结束（文件损坏或不完整）'
        except Exception as e:
            return True, f'PDF 读取失败: {type(e).__name__}'
        finally:
            logger.setLevel(old_level)
    except ImportError:
        return False, ''


def detect_broken_txt(file_path: Path) -> tuple[bool, str]:
    """检查 .txt 文件编码是否无法识别。"""
    if file_path.suffix.lower() != '.txt':
        return False, ''
    _TXT_ENCODINGS = ["utf-8", "gbk", "gb2312", "gb18030", "latin-1"]
    for enc in _TXT_ENCODINGS:
        try:
            content = file_path.read_text(encoding=enc)
            if content.strip():
                return False, ''
        except (UnicodeDecodeError, UnicodeError):
            continue
    return True, f'无法用常见编码 ({"/".join(_TXT_ENCODINGS)}) 解码'


def scan_compatibility(directory: Path) -> list[tuple[Path, str]]:
    """扫描目录，返回 (问题文件路径, 原因) 列表。"""
    problems = []
    for fp in sorted(directory.rglob('*')):
        if not fp.is_file():
            continue
        ext = fp.suffix.lower()
        if ext == '.docx':
            is_prob, reason = detect_docm(fp)
        elif ext == '.xlsx':
            is_prob, reason = detect_legacy_xls(fp)
        elif ext == '.pdf':
            is_prob, reason = detect_broken_pdf(fp)
        elif ext == '.txt':
            is_prob, reason = detect_broken_txt(fp)
        else:
            continue
        if is_prob:
            problems.append((fp, reason))
    return problems


def _convert_with_docling(source_path: Path, output_path: Path
                          ) -> Tuple[str, Optional[str], Optional[str]]:
    """
    用 Docling 转换文档（支持 docx/pdf）。

    注：不再在这里单独起子进程隔离 Docling 调用。上层的
    `convert_single_file_isolated()` 已经把"整个单文件转换"（包括这里的
    Docling 调用、失败后的原生降级路径等）打包放进一个独立子进程里跑，
    并配有真正的操作系统级超时+崩溃隔离；如果这里再单独起一层子进程，
    只会让每个 docx/pdf 文件多付一次 Python 解释器启动 + Docling 模型
    加载的开销，没有任何额外收益。
    返回 (status, error, warning)。

    ────────────────────────────────────────────────────────────────
    关于大页数 PDF 的说明
    ────────────────────────────────────────────────────────────────
    Docling 的 docling-parse C++ 后端在处理单个 PDF 时，preprocess
    阶段会按页累积内存，不会随页面处理完就释放（这是 Docling 自身
    已知的问题，参见其 GitHub issue #3345）。这个累积速度和每页的
    图片/内容密度直接相关，不能单纯用页数去预判"多少页以上才会出
    问题"——实测中一份 120 页、图片密集的硬件安装指导文档，在没有
    分批、单次 convert() 处理到第 117 页时就已经 std::bad_alloc，
    远早于社区案例里"数百页文本类 PDF 才会 OOM"的经验值。

    因此策略是：只要页数超过 PDF_BATCH_PAGE_SIZE 就分批，不设置更高
    的"大文件才分批"的门槛。分批时用 Docling 原生的 page_range 参数
    直接在原文件上取页，不用 pypdf 物理拆出临时文件再重新解析（后者
    多一次完整的读+写+重新解析开销）。每一批都是独立的 convert()
    调用，preprocess 阶段按页累积的内存会随每次新调用重新从零开始，
    从而把峰值内存压在"一批的量"而不是"全文档的量"上。

    另外，OCR（RapidOCR/onnxruntime）已经在 `_get_docling_converter()`
    里通过 `do_ocr = False` 显式关闭——知识库源文档是可直接抽取文本层
    的电子版文件，不需要识别图片里的文字；这也顺带避免了 OCR 推理时
    在大图片上分配内存失败导致整个子进程被系统强杀的问题
    （exitcode 3221225477 / 0xC0000005 access violation）。
    """
    ext = source_path.suffix.lower()
    if ext == ".pdf":
        return _convert_pdf_with_docling(source_path, output_path)

    try:
        converter = _get_docling_converter()
        result = converter.convert(str(source_path))
        md_content = result.document.export_to_markdown()
    except Exception as e:
        return ("error", f"Docling: {type(e).__name__}: {e}", None)

    if not md_content.strip():
        return ("empty", "Docling 转换内容为空", None)

    md_content = _clean_md_content(md_content)
    if not md_content.strip():
        return ("empty", "内容为空（移除了所有模板段落）", None)

    _ensure_output_dir(output_path)
    output_path.write_text(md_content, encoding="utf-8")
    warning = _check_md_size(output_path)

    return ("ok", None, warning)


# 每批处理的页数。之前认为"批次切得越小、convert()调用次数越多，
# 内存累积越明显"是针对 OCR（RapidOCR/onnxruntime）反复加载推理模型
# 这一具体场景的结论——现在已经关闭 OCR（do_ocr=False），这个顾虑
# 不再成立。
#
# 真正观察到的崩溃现场是 Docling 的 docling-parse C++ 后端在
# "preprocess"阶段按页处理时的内存累积，这一累积发生在单次 convert()
# 调用内部、按页递增，和该 PDF 每页图片/内容的密度直接相关，
# 不是单纯"页数多少"能预判的——这份 120 页的硬件安装指导文档在没有
# 分批、单次 convert() 里处理到第 117 页时就已经 std::bad_alloc，
# 说明"几百页才会出问题"的经验阈值对图片密集型文档完全不适用。
#
# 因此改为：只要页数超过下面这个批次大小，就用 Docling 原生的
# page_range 参数分批调用——每一批都是一次新的 convert() 调用，
# preprocess 阶段的累积状态会随之重置，从而把峰值内存压在
# "一批的量"而不是"全文档的量"上。批次大小不必也不该设得很大，
# 40 页是留了余量的保守值（远小于上面这份文件实际撑到的 117 页）。
PDF_BATCH_PAGE_SIZE = 40


def _convert_pdf_with_docling(source_path: Path, output_path: Path
                              ) -> Tuple[str, Optional[str], Optional[str]]:
    """
    PDF 专用转换路径。

    页数不超过 PDF_BATCH_PAGE_SIZE 时直接一次性 convert()。
    超过则用 Docling 原生的 page_range 参数按 PDF_BATCH_PAGE_SIZE
    分批调用——直接在原文件路径上取页，不会用 pypdf 另外拆出临时
    PDF 再重新解析，每一批都是独立的 convert() 调用，preprocess
    阶段在 C++ 后端里按页累积的内存会随每次新调用重置。
    """
    try:
        from pypdf import PdfReader
        total_pages = len(PdfReader(str(source_path)).pages)
    except Exception:
        # 拿不到页数（比如文件本身就有问题）时，退回一次性转换，
        # 把真正的错误交给下面的 convert() 去暴露，而不是在这里
        # 提前吞掉、给出一个无关的"读取页数失败"。
        total_pages = 0

    if total_pages <= PDF_BATCH_PAGE_SIZE:
        return _docling_convert_once(source_path, output_path)

    return _docling_convert_batched(source_path, output_path, total_pages)


def _docling_convert_once(source_path: Path, output_path: Path
                          ) -> Tuple[str, Optional[str], Optional[str]]:
    """单次整份 convert()，用于页数未超过阈值的 PDF（以及所有 docx）。"""
    try:
        converter = _get_docling_converter()
        result = converter.convert(str(source_path))
        md_content = result.document.export_to_markdown()
    except Exception as e:
        return ("error", f"Docling: {type(e).__name__}: {e}", None)

    if not md_content.strip():
        return ("empty", "Docling 转换内容为空", None)

    md_content = _clean_md_content(md_content)
    if not md_content.strip():
        return ("empty", "内容为空（移除了所有模板段落）", None)

    _ensure_output_dir(output_path)
    output_path.write_text(md_content, encoding="utf-8")
    warning = _check_md_size(output_path)

    return ("ok", None, warning)


def _docling_convert_batched(source_path: Path, output_path: Path,
                              total_pages: int
                             ) -> Tuple[str, Optional[str], Optional[str]]:
    """
    仅用于页数超过 PDF_BATCH_PAGE_SIZE 的 PDF。
    用 Docling 原生 page_range 参数按 PDF_BATCH_PAGE_SIZE 分批取页，
    直接在原文件上操作，不额外拆临时文件。
    """
    import gc

    converter = _get_docling_converter()
    md_parts: List[str] = []
    batch_count = (total_pages + PDF_BATCH_PAGE_SIZE - 1) // PDF_BATCH_PAGE_SIZE

    for i in range(batch_count):
        start = i * PDF_BATCH_PAGE_SIZE + 1          # Docling page_range 从 1 开始
        end = min((i + 1) * PDF_BATCH_PAGE_SIZE, total_pages)
        try:
            result = converter.convert(str(source_path), page_range=(start, end))
            part = result.document.export_to_markdown()
        except Exception as e:
            return ("error",
                     f"Docling（第 {start}-{end} 页，共 {total_pages} 页）: "
                     f"{type(e).__name__}: {e}", None)

        if part.strip():
            md_parts.append(part)

        # 及时丢弃这一批的中间结果，减少 convert() 反复调用带来的
        # 内存累积（Docling 底层 C++ 后端的累积无法靠这一步彻底消除，
        # 但能把 Python 侧能回收的部分尽量回收掉）。
        del result
        gc.collect()

    if not md_parts:
        return ("empty", "Docling 转换内容为空", None)

    md_content = "\n\n".join(md_parts).strip()
    md_content = _clean_md_content(md_content)
    if not md_content.strip():
        return ("empty", "内容为空（移除了所有模板段落）", None)

    _ensure_output_dir(output_path)
    output_path.write_text(md_content, encoding="utf-8")
    warning = _check_md_size(output_path)

    return ("ok", None, warning)


# ============================================================
# MD 文件大小预警
# ============================================================

def _check_md_size(md_path: Path) -> Optional[str]:
    """MD 文件超过阈值时返回预警信息。"""
    if md_path.exists():
        size_kb = md_path.stat().st_size / 1024
        if size_kb > MAX_MD_FILE_SIZE_KB:
            return f"MD 文件过大 ({size_kb:.0f}KB)，可能影响入库性能"
    return None


# ============================================================
# MD 内容清洁：移除封面/目录/版权/版本记录/作者信息等
# ============================================================

# 需要移除的单行模式（匹配即移除该行）
_BOILERPLATE_LINE_PATTERNS = [
    re.compile(r'^#+\s*(前\s*言|引言|概述|背景|前\s*言\s*$)'),
    re.compile(r'^#+\s*(目\s*录|目录|Contents)'),
    re.compile(r'(版权|著作权|著作权声明|版权声明|©\s*\d{4})'),
    re.compile(r'(All\s+rights?\s+reserved)', re.IGNORECASE),
    re.compile(r'(Confidential|机密|秘密|绝密)'),
    re.compile(r'(版本\s*[：:]|版\s*本\s*[：:]|版\s*本\s*号)'),
    re.compile(r'(修订记录|修\s*订\s*记\s*录|变更记录|变更历史|修订历史|文档变更)'),
    re.compile(r'^(作者|编写[：:]?|编制[：:]?|审核[：:]?|批准[：:]?|校对[：:]?|会审[：:]?|评审[：:]?|起草[：:]?|复核[：:]?)'),
    re.compile(r'^(第\s*\d+\s*页|Page\s+\d+|—\s*\d+\s*—|-\s*\d+\s*-)$'),
    re.compile(r'^中兴通讯|^ZTE\s*CORPORATION|^ZTE\s*中兴'),
    re.compile(r'^(技术文件|技术手册|产品文档|产品手册|用户手册|操作指南)'),
    re.compile(r'^文档版本\s*\d'),
    re.compile(r'^[\u2460-\u2473①-⑳]'),  # 带圈数字（常用于版权脚注）
]

# 目录检测
_TOC_ENTRY = re.compile(
    r'^\d+(\.\d+)*\s*[\u4e00-\u9fff][\u4e00-\u9fff\w]*[\s.…·]+\d+\s*$'
)
_TOC_NUM_LINE = re.compile(r'^[\d一二三四五六七八九十]+[.、．]\s*\S{1,40}$')
_TOC_PURE_DOTS = re.compile(r'^[\s.…·]+$')
_TOC_HEADING = re.compile(r'^#+\s*[目\t ]*[录録]\s*$|^[#\s]*目录|^[#\s]*Contents')


def _is_toc_section(lines: list[str], start: int, max_lookahead: int = 50) -> int:
    """检测目录段落，返回段落结束行号。"""
    if not _TOC_HEADING.match(lines[start].strip()):
        return start

    end = min(start + max_lookahead, len(lines))
    toc_count = 1
    for i in range(start + 1, end):
        raw = lines[i].strip()
        if not raw:
            toc_count += 1
            continue
        if raw.startswith('#'):
            continue
        if _TOC_PURE_DOTS.match(raw):
            toc_count += 1
            continue
        if _TOC_ENTRY.match(raw) or _TOC_NUM_LINE.match(raw):
            toc_count += 1
            continue
        break
    return start + toc_count if toc_count >= 4 else start


# 封面检测：常见封面行（短行，无编号无标题）
_COVER_LINE = re.compile(r'^[\u4e00-\u9fff\w\s]{1,30}$')
_COVER_END_MARKER = re.compile(r'^#{1,6}\s|^第[一二三四五六七八九十\d]+[章节篇]')


def _is_junk_table_row(cells: list[str]) -> bool:
    """判断表格行是否全是空单元格或纯零值（垃圾行）。"""
    for c in cells:
        s = c.strip()
        if s and s != '0':
            return False
    return True


def _compact_table_block(lines: list[str]) -> list[str]:
    """压缩表格块：去掉尾部空单元格、移除纯空/纯零行、自动计算实际列数。"""
    # 找出所有分隔线行（只检测实际表格中的分隔线 = 整行仅含 --- 和 |）
    # 并记录分隔线的列数
    sep_indices = set()
    sep_col_count = 0
    for i, line in enumerate(lines):
        cells = line.split('|')
        has_sep = sum(1 for c in cells if c.strip().replace('-', '').strip() == '' and '---' in c)
        if has_sep >= 2:
            sep_indices.add(i)
            if has_sep > sep_col_count:
                sep_col_count = has_sep

    if not sep_indices:
        # 没有分隔线：检查是否全是纯空/纯零行 → 移除
        all_junk = True
        for line in lines:
            cells = line.split('|')
            if len(cells) >= 3:
                content = [c.strip() for c in cells[1:-1]]
                if not _is_junk_table_row(content):
                    all_junk = False
                    break
        if all_junk:
            return []
        # 有内容行但仍可能有尾部空单元格，独立修剪
        out = []
        for line in lines:
            cells = line.split('|')
            if len(cells) < 3:
                out.append(line)
                continue
            cc = [c.strip() for c in cells[1:-1]]
            # 去掉尾部空单元格
            while cc and not cc[-1]:
                cc.pop()
            # 去掉尾部纯零单元格
            while cc and all(c == '0' for c in cc[-1:]):
                cc.pop()
            # 去掉前导空单元格
            while cc and not cc[0]:
                cc.pop(0)
            if _is_junk_table_row(cc):
                continue
            out.append('| ' + ' | '.join(cc) + ' |')
        return out

    if sep_col_count <= 1:
        return lines  # 单列分隔线，不处理

    # 处理每行：先去尾部空单元格，记录实际使用列数
    processed = []
    for i, line in enumerate(lines):
        cells = line.split('|')
        if len(cells) < 3:
            processed.append((i, cells, 0, False))
            continue

        content_cells = cells[1:-1]

        if i in sep_indices:
            col_count = sum(1 for c in content_cells if '---' in c)
            processed.append((i, ['---'] * col_count, col_count, True))
            continue

        # 数据行：去掉尾部空单元格
        while content_cells and not content_cells[-1].strip():
            content_cells.pop()
        while content_cells and all(c.strip() == '0' for c in content_cells[-1:]):
            content_cells.pop()

        processed.append((i, content_cells, len(content_cells), False))

    # 确定表格实际列数 = 所有数据行的最大列数（分隔线行除外）
    real_cols = max(
        (cnt for _, _, cnt, is_sep in processed if not is_sep and cnt > 0),
        default=0
    )
    if real_cols == 0:
        return []  # 纯空表（只有分隔线和空行），全部移除
    if real_cols == 1:
        return lines  # 单列表格，不处理

    out = []
    for idx, content_cells, _, is_sep in processed:
        if is_sep:
            sep_part = ['---'] * real_cols  # 强制补齐到实际列数
            out.append('|' + '|'.join(sep_part) + '|')
            continue

        if _is_junk_table_row(content_cells):
            continue

        cells = content_cells[:real_cols]
        reconstructed = '| ' + ' | '.join(c.strip() for c in cells) + ' |'
        out.append(reconstructed)

    return out


def _clean_md_content(content: str) -> str:
    """
    移除 Markdown 中的封面/目录/版权/版本记录/作者信息等模板化段落，
    以及表格中的空单元格/零值占位等冗余内容。
    策略：
      1. 先检测 TOC 段落（基于段落特征），标注范围
      2. 再移除匹配的单行模板模式
      3. 检测并移除文件开头的封面段
      4. 压缩表格：去掉尾部空单元格、移除纯空/纯零行
    """
    if not content.strip():
        return content

    lines = content.split('\n')
    n = len(lines)
    keep = [True] * n

    # Pass 1: 目录段落
    i = 0
    while i < n:
        toc_end = _is_toc_section(lines, i)
        if toc_end > i:
            for j in range(i, toc_end):
                keep[j] = False
            i = toc_end
            continue
        i += 1

    # Pass 2: 单行模板模式
    for i, line in enumerate(lines):
        if not keep[i]:
            continue
        stripped = line.strip()
        if any(p.search(stripped) for p in _BOILERPLATE_LINE_PATTERNS):
            keep[i] = False

    # Pass 3: 封面段
    first_heading_idx = -1
    for i, line in enumerate(lines):
        if not keep[i]:
            continue
        if re.match(r'^#{1,6}\s', line.strip()):
            first_heading_idx = i
            break

    if first_heading_idx > 0:
        # 第一个标题前的行 -> 如果全是短封面行 -> 移除
        pre_kept = [lines[j] for j in range(first_heading_idx) if keep[j] and lines[j].strip()]
        if pre_kept and all(len(l.strip()) < 40 for l in pre_kept):
            for j in range(first_heading_idx):
                keep[j] = False

    # Pass 4: 表格压缩（对每个表格块单独处理）
    result_lines = []
    i = 0
    while i < n:
        if not keep[i]:
            result_lines.append(lines[i])
            i += 1
            continue

        # 检测表格块：只要包含 | 且在分隔线行或空/零值行范围内的连续行
        raw = lines[i].strip()
        if '|' in raw:
            block = []
            while i < n and keep[i]:
                cur = lines[i].strip()
                if '|' not in cur:
                    break
                # 规范化：确保有前导 | 和尾随 | 方便处理
                if not cur.startswith('|'):
                    cur = '|' + cur
                if not cur.endswith('|'):
                    cur = cur + '|'
                block.append(cur)
                i += 1
            if block:
                compressed = _compact_table_block(block)
                result_lines.extend(compressed)
            continue

        result_lines.append(lines[i])
        i += 1

    cleaned = '\n'.join(result_lines)

    # ═══ 兜底核弹清洁 ═══
    #  清理 _compact_table_block 可能遗漏的表格垃圾
    #  1) 纯空行 & 纯 --- 行 → 整行删除
    cleaned = re.sub(r'^\|(?:\s*\|\s*)*\|\s*$', '', cleaned, flags=re.M)
    #  2) 空行边界合并（删除空表后可能留的多余空行）
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    #  3) 每个表格行尾部空单元格截断：| a | b |   |   | → | a | b |
    #     去掉行尾连续的 | (空格) 模式
    cleaned = re.sub(r'(?<=\|)(?:\s*\|\s*)*$', '', cleaned, flags=re.M)
    #  4) 去掉前导空单元格：|   |   | a | b | → | a | b |
    cleaned = re.sub(r'^\|(?:\s*\|\s*)*\s*(?=\S)', '| ', cleaned, flags=re.M)
    #  5) 清理残留的孤立空单元格（行中连续多个空格|）  
    cleaned = re.sub(r'\|\s*\|\s*(?=\|)', '| ', cleaned)
    #  6) 移除 Docling 残留的 HTML 注释（如 <!-- image -->、<!-- 表1-1 指标集示例 --> 等）
    cleaned = re.sub(r'<!--.*?-->', '', cleaned, flags=re.DOTALL)

    cleaned = re.sub(r'\n{4,}', '\n\n\n', cleaned)
    return cleaned.strip()


# ============================================================
# 原生降级解析器
# ============================================================

def _docx_fallback(source_path: Path, output_path: Path
                   ) -> Tuple[str, Optional[str]]:
    """python-docx 降级方案（Docling 不可用时）。"""
    try:
        from docx import Document
        doc = Document(str(source_path))
        return _docx_to_md(doc, source_path, output_path)
    except ValueError as e:
        if "not a Word file" in str(e):
            return _docx_fallback_zip(source_path, output_path)
        return ("error", f"{type(e).__name__}: {e}")
    except PackageNotFoundError:
        return ("error", "不是有效的 .docx 文件（ZIP 包无法解析）")
    except KeyError as e:
        if "'NULL'" in str(e) or "NULL" in str(e):
            return _docx_fallback_zip(source_path, output_path)
        return ("error", f"{type(e).__name__}: {e}")
    except Exception as e:
        e_str = f"{type(e).__name__}: {e}"
        e_low = e_str.lower()
        if "no relationship" in e_low and "officedocument" in e_low \
           or "'null'" in e_low or "xmlsyntaxerror" in e_low \
           or ("xml" in e_low and "expected" in e_low):
            return _docx_fallback_zip(source_path, output_path)
        return ("error", e_str)


def _para_to_md_line(para) -> Optional[str]:
    text = para.text.strip()
    if not text:
        return None
    style_name = para.style.name.lower() if para.style else ""
    if "heading 1" in style_name or "title" in style_name:
        return f"# {text}"
    if "heading 2" in style_name:
        return f"## {text}"
    if "heading 3" in style_name:
        return f"### {text}"
    if "heading" in style_name:
        level = 4
        for s in ("heading 4", "heading 5", "heading 6"):
            if s in style_name:
                level = int(s[-1])
                break
        return f"{'#' * level} {text}"
    return text


def _table_to_md_lines(table) -> List[str]:
    lines = [""]
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        lines.append("| " + " | ".join(cells) + " |")
    lines.append("")
    return lines


def _docx_to_md(doc, source_path: Path, output_path: Path) -> Tuple[str, Optional[str]]:
    # 按文档原始顺序遍历正文元素（段落/表格交错出现），而不是先把所有
    # 段落收集完再统一把所有表格追加到末尾——后者会把"表格前后各有一段
    # 说明文字"的排版完全打乱，表格永远漂到文档最后，语义关联丢失。
    # python-docx 没有直接给出"正文元素按序遍历"的 API，但 Paragraph 和
    # Table 对象都能拿到各自的底层 XML 元素（`_p` / `_tbl`），可以反查
    # 它在 body 里的原始位置。
    body = doc.element.body
    para_by_elem = {p._p: p for p in doc.paragraphs}
    table_by_elem = {t._tbl: t for t in doc.tables}

    md_lines = []
    for child in body.iterchildren():
        if child in para_by_elem:
            line = _para_to_md_line(para_by_elem[child])
            if line is not None:
                md_lines.append(line)
        elif child in table_by_elem:
            md_lines.extend(_table_to_md_lines(table_by_elem[child]))

    md_content = "\n\n".join(md_lines).strip()
    if not md_content:
        return ("empty", "文档内容为空")

    md_content = _clean_md_content(md_content)
    if not md_content:
        return ("empty", "文档内容为空（移除了所有模板段落）")

    _ensure_output_dir(output_path)
    output_path.write_text(md_content, encoding="utf-8")
    if not is_file_readable(output_path):
        return ("garbled", "转换后内容疑似乱码")
    return ("ok", None)


def _docx_fallback_zip(source_path: Path, output_path: Path) -> Tuple[str, Optional[str]]:
    try:
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(str(source_path)) as z:
            for p in ["word/document.xml", "Word/document.xml"]:
                try:
                    xml_content = z.read(p)
                    break
                except KeyError:
                    continue
            else:
                return ("error", "ZIP 中未找到 word/document.xml")
        root = ET.fromstring(xml_content)
        ns_w = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        md_lines = []
        for para in root.iter(f"{{{ns_w}}}p"):
            texts = [t.text for t in para.iter(f"{{{ns_w}}}t") if t.text]
            line = "".join(texts).strip()
            if line:
                md_lines.append(line)
        md_content = "\n\n".join(md_lines).strip()
        if not md_content:
            return ("empty", "ZIP 回退提取内容为空")
        md_content = _clean_md_content(md_content)
        if not md_content:
            return ("empty", "ZIP 回退提取内容为空（移除模板后）")
        _ensure_output_dir(output_path)
        output_path.write_text(md_content, encoding="utf-8")
        return ("ok", None)
    except Exception as e:
        return ("error", f"ZIP 回退提取失败: {type(e).__name__}: {e}")


def _pdf_fallback(source_path: Path, output_path: Path
                  ) -> Tuple[str, Optional[str]]:
    """pypdf 降级方案（Docling 不可用时）。"""
    try:
        from pypdf import PdfReader, errors as pypdf_errors
        try:
            reader = PdfReader(str(source_path))
        except pypdf_errors.PdfStreamError:
            return ("error", "PDF 流意外结束（文件可能损坏或不完整）")
        except Exception as e:
            return ("error", f"PDF 读取失败: {type(e).__name__}: {e}")

        md_lines = []
        for i, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text()
            except Exception:
                text = ""
            if text and text.strip():
                md_lines.append(f"<!-- Page {i} -->\n\n{text.strip()}")

        md_content = "\n\n".join(md_lines).strip()
        if not md_content:
            return ("empty", "PDF 内容为空（可能是扫描件）")
        md_content = _clean_md_content(md_content)
        if not md_content.strip():
            return ("empty", "PDF 内容为空（移除了所有模板段落）")
        _ensure_output_dir(output_path)
        output_path.write_text(md_content, encoding="utf-8")
        if not is_file_readable(output_path):
            return ("garbled", "转换后内容疑似乱码")
        return ("ok", None)
    except ImportError:
        return ("error", "缺少 pypdf 库")
    except Exception as e:
        return ("error", f"PyPDF {type(e).__name__}: {e}")


def _xlsx_fallback(source_path: Path, output_path: Path
                   ) -> Tuple[str, Optional[str]]:
    """openpyxl 降级方案（Docling 不可用时）。"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(source_path), read_only=True, data_only=True)
        md_lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            md_lines.append(f"## {sheet_name}")
            md_lines.append("")
            row_iter = ws.iter_rows(values_only=True)
            try:
                header_row = next(row_iter)
            except StopIteration:
                continue
            headers = [str(c) if c is not None else "" for c in header_row]
            md_lines.append("| " + " | ".join(headers) + " |")
            md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
            for row in row_iter:
                if _is_junk_xlsx_row(row):
                    continue
                cells = [str(c) if c is not None else "" for c in row]
                md_lines.append("| " + " | ".join(cells) + " |")
            md_lines.append("")
        wb.close()
        md_content = "\n".join(md_lines).strip()
        if not md_content:
            return ("empty", "Excel 内容为空")
        md_content = _clean_md_content(md_content)
        if not md_content.strip():
            return ("empty", "Excel 内容为空（移除了所有模板段落）")
        _ensure_output_dir(output_path)
        output_path.write_text(md_content, encoding="utf-8")
        return ("ok", None)
    except ImportError:
        return ("error", "缺少 openpyxl 库")
    except Exception as e:
        return ("error", f"{type(e).__name__}: {e}")


def _pptx_fallback(source_path: Path, output_path: Path
                   ) -> Tuple[str, Optional[str]]:
    """python-pptx 降级方案（Docling 不可用时）。"""
    try:
        from pptx import Presentation
        prs = Presentation(str(source_path))
        md_lines = []
        for i, slide in enumerate(prs.slides, 1):
            md_lines.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            md_lines.append(text)
            md_lines.append("")
        md_content = "\n".join(md_lines).strip()
        if not md_content:
            return ("empty", "幻灯片内容为空")
        md_content = _clean_md_content(md_content)
        if not md_content.strip():
            return ("empty", "幻灯片内容为空（移除了所有模板段落）")
        _ensure_output_dir(output_path)
        output_path.write_text(md_content, encoding="utf-8")
        return ("ok", None)
    except ImportError:
        return ("error", "缺少 python-pptx 库")
    except Exception as e:
        return ("error", f"{type(e).__name__}: {e}")


# ============================================================
# xlsx 空行/零值行过滤（供降级用）
# ============================================================

def _is_junk_cell(value) -> bool:
    if value is None:
        return True
    if isinstance(value, (int, float)):
        return value == 0 or value == 0.0
    s = str(value).strip()
    return s == "" or s == "0" or s == "0.0"


def _is_junk_xlsx_row(row: tuple) -> bool:
    if not row:
        return True
    return all(_is_junk_cell(cell) for cell in row)


# ============================================================
# Markdown (.md) — 直接复制
# ============================================================

def _convert_md(source_path: Path, output_path: Path) -> Tuple[str, Optional[str], Optional[str]]:
    try:
        content = source_path.read_text(encoding="utf-8")
        if not content.strip():
            return ("empty", "文件内容为空", None)
        if not is_file_readable(source_path):
            return ("garbled", "源文件内容疑似乱码", None)
        content = _clean_md_content(content)
        if not content.strip():
            return ("empty", "文件内容为空（移除了所有模板段落）", None)
        _ensure_output_dir(output_path)
        output_path.write_text(content, encoding="utf-8")
        warning = _check_md_size(output_path)
        return ("ok", None, warning)
    except Exception as e:
        return ("error", f"{type(e).__name__}: {e}", None)


# ============================================================
# 代码/配置类文件 (.sql/.yaml/.yml/.json/.ini/.conf/.toml)
# ============================================================
# 这类文件不走 _clean_md_content()：那套清洗是给"技术文档"调的（剥封面、
# 剥版本记录行、压缩表格……），SQL/YAML 里类似 "-- 版本: 1.0" 的注释行、
# YAML 里的 "author: xxx" 字段，会被误当成模板行整行删掉；纯文本压缩表格
# 那一段碰到 SQL 里出现 "|" 字符（位运算/管道）也可能误触发。
# 直接原样保留内容，用 fenced code block 包一层方便阅读，并在文件首行
# 写入 SOURCE_EXT_MARKER_PREFIX 标记，供 ingest.py 分块阶段识别原始类型、
# 选择语义切分策略（按 SQL 语句/YAML 顶层 key/INI section 切，而不是按
# 字符数硬切，避免把一条 CREATE TABLE 语句从中间切断）。

_CODE_LANG_TAG = {
    ".sql": "sql", ".yaml": "yaml", ".yml": "yaml",
    ".json": "json", ".ini": "ini", ".conf": "ini", ".toml": "toml",
}


def _convert_code_config(source_path: Path, output_path: Path) -> Tuple[str, Optional[str], Optional[str]]:
    ext = source_path.suffix.lower()
    for enc in _TXT_ENCODINGS:
        try:
            content = source_path.read_text(encoding=enc)
            if content.strip():
                break
        except (UnicodeDecodeError, UnicodeError):
            continue
    else:
        return ("error", f"无法用 {', '.join(_TXT_ENCODINGS)} 解码文件内容", None)

    if not content.strip():
        return ("empty", "文件内容为空", None)
    if not is_file_readable(source_path):
        return ("garbled", "源文件内容疑似乱码", None)

    lang = _CODE_LANG_TAG.get(ext, "")
    marker = f"{SOURCE_EXT_MARKER_PREFIX}{ext} -->"
    md_content = f"{marker}\n\n```{lang}\n{content.rstrip()}\n```\n"

    _ensure_output_dir(output_path)
    output_path.write_text(md_content, encoding="utf-8")
    warning = _check_md_size(output_path)
    return ("ok", None, warning)


# ============================================================
# 纯文本 (.txt) — 支持编码回退
# ============================================================

_TXT_ENCODINGS = ["utf-8", "gbk", "gb2312", "gb18030", "latin-1"]


def _convert_txt(source_path: Path, output_path: Path) -> Tuple[str, Optional[str], Optional[str]]:
    for enc in _TXT_ENCODINGS:
        try:
            content = source_path.read_text(encoding=enc)
            if content.strip():
                break
        except (UnicodeDecodeError, UnicodeError):
            continue
    else:
        return ("error", f"无法用 {', '.join(_TXT_ENCODINGS)} 解码文件内容", None)
    if not content.strip():
        return ("empty", "文件内容为空", None)
    content = _clean_md_content(content)
    if not content.strip():
        return ("empty", "文件内容为空（移除了所有模板段落）", None)
    _ensure_output_dir(output_path)
    output_path.write_text(content, encoding="utf-8")
    warning = _check_md_size(output_path)
    return ("ok", None, warning)


# ============================================================
# 原生降级调度器
# ============================================================

_FALLBACK_REGISTRY = {
    ".docx": _docx_fallback,
    ".pdf":  _pdf_fallback,
    ".xlsx": _xlsx_fallback,
    ".pptx": _pptx_fallback,
}


def convert_single_file(source_path: Path) -> dict:
    """
    转换单个源文件为 Markdown。

    Parameters
    ----------
    source_path : Path
        源文件绝对路径。

    Returns
    -------
    dict: {
        "rel_path": ...,
        "status": "ok" | "garbled" | "empty" | "error" | "skip",
        "md_path": ... | None,
        "error": ... | None,
        "warning": ... | None,
        "sha256": ...,
        "size": ...,
        "mtime": ...,
    }
    """
    rel_path = _get_rel_path(source_path)
    ext = source_path.suffix.lower()
    result = {
        "rel_path": rel_path,
        "status": "error",
        "md_path": None,
        "error": None,
        "warning": None,
        "sha256": "",
        "size": 0,
        "mtime": "",
    }

    try:
        stat = source_path.stat()
        result["size"] = stat.st_size
        result["sha256"] = compute_sha256(source_path)

        from datetime import datetime, timezone
        dt = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)
        result["mtime"] = dt.isoformat(timespec="seconds")

        if ext not in SUPPORTED_EXTENSIONS:
            result["status"] = "skip"
            result["error"] = f"不支持格式 {ext}，忽略: {rel_path}"
            return result

        # 检测伪装为 .docx 的宏文档 (.docm) / 老版 .doc 误标；
        # 检测伪装为 .xlsx 的老版 .xls 误标。
        if ext == ".docx":
            is_problem, reason = detect_docm(source_path)
            if is_problem:
                if "宏文档" in reason:
                    # 宏文档（.docm）不能自动"另存为 .docx"——那样会把里面的
                    # VBA 代码直接丢掉，是否需要保留宏功能得让人来判断，
                    # 这一步不自动化，维持原来"跳过 + 提示手工处理"的行为。
                    result["status"] = "skip"
                    result["error"] = f"宏文档 (.docm)，需用 Word 另存为 .docx: {rel_path}"
                    return result
                # 大概率是老版 .doc 误标成了 .docx，尝试用 WPS/Office
                # 自动另存为真正的 .docx；失败了再退回手工处理。
                ok, msg = _convert_legacy_office_file(source_path, "word")
                if not ok:
                    result["status"] = "skip"
                    result["error"] = (f"{reason}，自动转换失败（{msg}），"
                                        f"需手工用 WPS/Word 另存为 .docx: {rel_path}")
                    return result
                result["warning"] = f"自动转换: {msg}"
                # 文件内容已经被原地替换，重新计算 size/sha256/mtime，
                # 后面转换流程按新内容走。
                stat = source_path.stat()
                result["size"] = stat.st_size
                result["sha256"] = compute_sha256(source_path)
                dt = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)
                result["mtime"] = dt.isoformat(timespec="seconds")

        elif ext == ".xlsx":
            is_problem, reason = detect_legacy_xls(source_path)
            if is_problem:
                ok, msg = _convert_legacy_office_file(source_path, "excel")
                if not ok:
                    result["status"] = "skip"
                    result["error"] = (f"{reason}，自动转换失败（{msg}），"
                                        f"需手工用 WPS/Excel 另存为 .xlsx: {rel_path}")
                    return result
                result["warning"] = f"自动转换: {msg}"
                stat = source_path.stat()
                result["size"] = stat.st_size
                result["sha256"] = compute_sha256(source_path)
                dt = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)
                result["mtime"] = dt.isoformat(timespec="seconds")

        output_path = _get_output_md_path(source_path)

        # ── 策略：Docling 主力(docx/pdf) → 原生降级(xlsx/pptx) → 直接复制(md/txt) ──
        if ext in _DOCLING_SUPPORTED:
            # 主力：Docling
            status, error, warning = _convert_with_docling(source_path, output_path)
            if status == "ok":
                result["status"] = status
                result["error"] = error
                result["warning"] = warning
            else:
                # Docling 失败 → 降级到原生解析器
                fallback = _FALLBACK_REGISTRY.get(ext)
                if fallback:
                    status_fb, error_fb = fallback(source_path, output_path)
                    if status_fb == "ok":
                        result["status"] = status_fb
                        result["error"] = None
                        result["warning"] = _check_md_size(output_path)
                    else:
                        result["status"] = status_fb
                        result["error"] = error_fb
                else:
                    result["status"] = status
                    result["error"] = error
        elif ext in _FALLBACK_REGISTRY:
            # xlsx/pptx：直接使用原生解析器
            fallback = _FALLBACK_REGISTRY[ext]
            status_fb, error_fb = fallback(source_path, output_path)
            result["status"] = status_fb
            result["error"] = error_fb
            if status_fb == "ok":
                result["warning"] = _check_md_size(output_path)
        elif ext == ".md":
            status, error, warning = _convert_md(source_path, output_path)
            result["status"] = status
            result["error"] = error
            result["warning"] = warning
        elif ext == ".txt":
            status, error, warning = _convert_txt(source_path, output_path)
            result["status"] = status
            result["error"] = error
            result["warning"] = warning
        elif ext in CODE_CONFIG_EXTENSIONS:
            status, error, warning = _convert_code_config(source_path, output_path)
            result["status"] = status
            result["error"] = error
            result["warning"] = warning

        if result["status"] == "ok":
            result["md_path"] = str(output_path.with_suffix(".md").absolute())

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"意外异常: {type(e).__name__}: {e}\n{traceback.format_exc()}"

    return result


# 通用单文件转换子进程 Worker 路径
_CONVERT_WORKER_PATH = Path(__file__).parent / "convert_worker.py"


def convert_single_file_isolated(source_path: Path, timeout: int = CONVERT_TIMEOUT,
                                 start_callback=None) -> dict:
    """
    在独立子进程里转换单个文件——这是真正的操作系统级超时 + 崩溃隔离。

    相比直接在 ThreadPoolExecutor 的线程里调用 convert_single_file()：
      - 超时保护是真的：子进程超时后会被 subprocess.run(timeout=...)
        用 SIGKILL 强制杀掉，对应的线程池 worker 线程能立刻返回、空出来
        去处理下一个排队的文件。不会再出现"某个文件在原生解析库里卡死，
        线程无法从外部打断，线程池被逐渐耗尽，as_completed() 因为没有
        任何 future 完成而无限期阻塞、界面停止刷新、Ctrl+C 后重跑仍卡在
        同一个文件"的情况（这正是本函数要修复的问题——旧版把超时判断
        写在 as_completed() 之后，只有 future 真正完成时才检查得到，
        对于卡死在线程里、永远不会完成的任务完全不起作用）。
      - 崩溃隔离：任何原生扩展（Docling / python-docx / openpyxl / ...）
        如果段错误，只有这一个子进程死掉，主流水线不受影响。

    代价：每个文件多一次 Python 解释器启动开销（通常 <1s）。docx/pdf 走
    Docling 时模型仍按文件重新加载——这是历史行为，不是本次改动引入的
    新开销（另见 README 的性能优化建议一节）。

    start_callback 在这里（真正开始跑这个文件的 worker 线程里）调用，
    而不是在提交任务的那一刻调用——线程池只有 max_workers 个并发名额，
    如果在"提交"时就打印"开始处理"，会把还在排队、根本没轮到执行的文件
    也一起打出来，跟实际进度完全对不上（之前的版本就是这个问题：384 个
    文件的"开始处理"几乎在同一秒内全部打印完，之后线程池才真正开始跑，
    看起来像是"全部同时卡住"，其实只是日志时机不对）。
    """
    import subprocess
    import tempfile
    import json as _json

    rel_path = _get_rel_path(source_path)

    if start_callback:
        try:
            start_callback(rel_path)
        except Exception:
            pass  # 日志/回调本身出问题，不能让它连累转换任务

    rf = tempfile.NamedTemporaryFile(mode="w", suffix=".json", delete=False)
    rf.close()
    result_file = rf.name

    result = {
        "rel_path": rel_path,
        "status": "error",
        "md_path": None,
        "error": None,
        "warning": None,
        "sha256": "",
        "size": 0,
        "mtime": "",
    }

    # 只有 docx/pdf 会真正走 Docling（拖内存的那部分）；xlsx/pptx/txt 等
    # 原生解析格式完全不需要排队等这个信号量，仍按 CONVERT_WORKERS 的
    # 并发数直接跑，不受影响。
    needs_docling_slot = source_path.suffix.lower() in _DOCLING_SUPPORTED
    if needs_docling_slot:
        _docling_semaphore.acquire()
    try:
        proc = subprocess.run(
            [sys.executable, str(_CONVERT_WORKER_PATH), str(source_path), result_file],
            capture_output=True,
            timeout=timeout,
        )
        if proc.returncode < 0:
            result["error"] = f"转换子进程被信号 {-proc.returncode} 杀死（可能段错误）"
            return result
        if proc.returncode != 0:
            stderr_tail = proc.stderr.decode("utf-8", "ignore")[-800:] if proc.stderr else ""
            # Windows 下常见的原生崩溃 exitcode（不会像 POSIX 那样变成负数），
            # 直接提示出来，比一串 exitcode 数字更容易看懂是什么问题。
            _WIN_CRASH_CODES = {
                3221225477: "0xC0000005 访问违规（常见于内存不足/耗尽时原生扩展崩溃）",
                3221225620: "0xC0000094 整数除零",
                3221225725: "0xC00000FD 栈溢出",
            }
            hint = _WIN_CRASH_CODES.get(proc.returncode)
            hint_str = f"，疑似 {hint}" if hint else ""
            result["error"] = f"转换子进程异常退出 (exitcode={proc.returncode}){hint_str} {stderr_tail}".strip()
            return result

        with open(result_file, encoding="utf-8") as f:
            loaded = _json.load(f)
        # 子进程正常返回但内部捕获了异常（比如 convert.py import 失败）时，
        # loaded 里可能只有 status/error，缺少 rel_path 等字段，这里补全。
        result.update(loaded)
        result.setdefault("rel_path", rel_path)
        return result

    except subprocess.TimeoutExpired:
        result["error"] = f"转换超时（>{timeout}s），已强制终止子进程并跳过该文件"
        return result
    except Exception as e:
        result["error"] = f"调用转换子进程失败: {type(e).__name__}: {e}"
        return result
    finally:
        try:
            os.unlink(result_file)
        except OSError:
            pass
        if needs_docling_slot:
            _docling_semaphore.release()


def convert_batch(file_paths: List[Path],
                  max_workers: int = CONVERT_WORKERS,
                  timeout: int = CONVERT_TIMEOUT,
                  progress_callback=None,
                  start_callback=None) -> List[dict]:
    """
    批量转换文件，支持 Ctrl+C 中断和单文件超时（真实生效，见
    convert_single_file_isolated 的说明）。

    Parameters
    ----------
    start_callback : callable, optional
        每个文件真正开始在 worker 里执行时调用一次（参数：相对路径），
        用于实时打印"正在处理 xxx"，便于在文件真的很慢/卡住时第一时间
        定位是哪个文件——注意这是"真正开始跑"才触发，不是"提交进队列"
        就触发，避免和线程池的实际并发数对不上。
    """
    results = []

    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        futures = {
            pool.submit(convert_single_file_isolated, fp, timeout, start_callback): fp
            for fp in file_paths
        }

        try:
            for future in as_completed(futures):
                fp = futures[future]
                try:
                    # 子进程内部已经用 subprocess.run(timeout=...) 强制
                    # 保证不会超过 `timeout` 秒，这里的 timeout 只是双重
                    # 保险（给子进程自身的收尾/清理留一点缓冲时间）。
                    result = future.result(timeout=timeout + 30)
                except TimeoutError:
                    rel = fp.relative_to(SOURCE_DIR).as_posix() if hasattr(fp, 'relative_to') else str(fp)
                    result = {
                        "rel_path": rel,
                        "status": "error",
                        "md_path": None,
                        "error": f"转换超时（>{timeout + 30}s），已跳过",
                        "warning": None,
                        "sha256": "",
                        "size": 0,
                        "mtime": "",
                    }
                    future.cancel()
                except Exception as e:
                    # worker 内部出现了 convert_single_file_isolated 自身
                    # 都没能捕获的异常（理论上不该发生，但防御性地兜住，
                    # 避免单个文件的意外问题带崩整个 as_completed 循环，
                    # 导致本该继续跑的其它文件也一起停摆）。
                    rel = fp.relative_to(SOURCE_DIR).as_posix() if hasattr(fp, 'relative_to') else str(fp)
                    result = {
                        "rel_path": rel,
                        "status": "error",
                        "md_path": None,
                        "error": f"worker 内部异常: {type(e).__name__}: {e}",
                        "warning": None,
                        "sha256": "",
                        "size": 0,
                        "mtime": "",
                    }
                results.append(result)
                if progress_callback:
                    try:
                        progress_callback(result)
                    except Exception as e:
                        # 同样道理：打印/记录进度这一步本身出错（比如
                        # Windows 控制台编码不支持某些字符），不能让它
                        # 把整个转换流程带崩、连累后面还没处理的文件。
                        print(f"[WARN] progress_callback 出错（已忽略，"
                              f"不影响转换结果）: {type(e).__name__}: {e}",
                              file=sys.stderr)
        except KeyboardInterrupt:
            for f in futures:
                f.cancel()
            raise
        finally:
            pool.shutdown(wait=False, cancel_futures=True)

    results.sort(key=lambda r: r["rel_path"])
    return results
